import os
import random
import tempfile
import base64
import subprocess
from fastapi import Form, Request, BackgroundTasks
from fastapi.responses import FileResponse, RedirectResponse
from pptx import Presentation


def get_random_number():
    """Получает случайное число до 6 знаков."""
    return random.randint(100000, 999999)


def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Конвертирует PPTX файл в PDF используя LibreOffice."""
    try:
        libreoffice_paths = [
            'libreoffice',
            '/usr/bin/libreoffice',
            '/usr/local/bin/libreoffice',
            '/snap/bin/libreoffice',
            '/opt/libreoffice/program/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice'
        ]
        libreoffice_cmd = None
        
        for path in libreoffice_paths:
            try:
                result = subprocess.run(
                    [path, '--version'], capture_output=True, timeout=5
                )
                if result.returncode == 0:
                    libreoffice_cmd = path
                    break
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue
        
        if not libreoffice_cmd:
            raise Exception(
                "LibreOffice не найден. Установите LibreOffice:\n"
                "Ubuntu/Debian: sudo apt-get install libreoffice\n"
                "CentOS/RHEL: sudo yum install libreoffice\n"
                "macOS: brew install --cask libreoffice"
            )
        
        cmd = [
            libreoffice_cmd,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', os.path.dirname(pdf_path),
            pptx_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
        
        if result.returncode != 0:
            raise Exception(f"LibreOffice error: {result.stderr}")
            
        
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        generated_pdf = os.path.join(os.path.dirname(pdf_path), f"{base_name}.pdf")
        
        if not os.path.exists(generated_pdf):
            raise Exception("PDF файл не был создан")
            
        if generated_pdf != pdf_path:
            os.rename(generated_pdf, pdf_path)
            
    except subprocess.TimeoutExpired:
        raise Exception("Превышено время ожидания конвертации")
    except Exception as e:
        raise Exception(f"Ошибка конвертации: {str(e)}")


def gen_cert_handler(
    request: Request,
    name: str | None = Form(None),
    price: str | None = Form(None)
):
    """Обработчик для генерации сертификатов"""
    try:
        name_value = name.strip() if name else ""
        price_value = price.strip() if price else ""
        
        if price_value:
            if len(price_value) > 6:
                raise ValueError("Цена слишком большая (максимум 999999)")
            
            try:
                price_int = int(price_value)
                if price_int <= 0:
                    raise ValueError("Цена должна быть больше 0")
            except ValueError:
                raise ValueError("Цена должна быть числом")
        
        template_path = os.path.join(
            os.path.dirname(__file__), 
            'Сертификат_шаблон.pptx'
        )
        
        if not os.path.exists(template_path):
            raise FileNotFoundError(
                f"Файл шаблона не найден: {template_path}"
            )
        
        prs = Presentation(template_path)
        
        serial_number = get_random_number()
        
        replacements = {
            'price': price_value,
            'name': name_value,
            'serial': str(serial_number),
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            if key in run.text:
                                run.text = run.text.replace(key, value)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_pptx:
            temp_pptx_path = temp_pptx.name
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            temp_pdf_path = temp_pdf.name
        
        prs.save(temp_pptx_path)
        
        convert_pptx_to_pdf(temp_pptx_path, temp_pdf_path)
        
        output_filename = "Сертификат.pdf"
        
        def cleanup_temp_files():
            try:
                os.unlink(temp_pptx_path)
                os.unlink(temp_pdf_path)
            except OSError:
                pass
        
        media_type = 'application/pdf'
        
        background_tasks = BackgroundTasks()
        background_tasks.add_task(cleanup_temp_files)
        
        return FileResponse(
            path=temp_pdf_path,
            filename=output_filename,
            media_type=media_type,
            background=background_tasks
        )
        
    except Exception as e:
        status = f"Ошибка генерации сертификата: {str(e)}"
        response = RedirectResponse(url="/gen_rit_cert", status_code=303)
        encoded_status = base64.b64encode(status.encode('utf-8')).decode('ascii')
        response.set_cookie("gen_cert_status", encoded_status, max_age=10)
        return response