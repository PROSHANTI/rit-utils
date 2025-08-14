import datetime
import locale
import tempfile
import os
import base64
from fastapi import Form, Request, BackgroundTasks
from fastapi.responses import FileResponse, RedirectResponse
from pptx import Presentation

try:
    locale.setlocale(locale.LC_ALL, 'ru_RU.UTF-8')
except locale.Error:
    try:
        locale.setlocale(locale.LC_ALL, 'ru_RU')
    except locale.Error:
        try:
            locale.setlocale(locale.LC_ALL, 'en_US.UTF-8')
        except locale.Error:
            try:
                locale.setlocale(locale.LC_ALL, 'C.UTF-8')
            except locale.Error:
                # Если все локали недоступны, используем дефолтную
                locale.setlocale(locale.LC_ALL, 'C')


def get_current_date():
    """Получить текущую дату и вернуть день, месяц и год."""
    now = datetime.datetime.now()
    try:
        month_name = now.strftime("%B")
    except (OSError, ValueError):
        month_mapping = {
            'January': 'января', 'February': 'февраля', 'March': 'марта',
            'April': 'апреля', 'May': 'мая', 'June': 'июня',
            'July': 'июля', 'August': 'августа', 'September': 'сентября',
            'October': 'октября', 'November': 'ноября', 'December': 'декабря'
        }
        month_name = month_mapping.get(now.strftime("%B"), now.strftime("%B"))
    
    return now.day, month_name, now.year


def doctor_form_handler(
    request: Request,
    doctor_1: str | None = Form(None),
    doctor_2: str | None = Form(None), 
    doctor_3: str | None = Form(None),
    doctor_4: str | None = Form(None),
    patient_1: str | None = Form(None),
    patient_2: str | None = Form(None),
    patient_3: str | None = Form(None),
    patient_4: str | None = Form(None),
    date: str | None = Form(None)
):
    try:
        day, month, year = get_current_date()
        
        if date and date.strip() and date.isdigit():
            day = int(date)
        
        template_path = os.path.join(
            os.path.dirname(__file__), 
            'Бланк Врача.pptx'
        )
        
        if not os.path.exists(template_path):
            raise FileNotFoundError(
                f"Файл шаблона не найден: {template_path}"
            )
        
        prs = Presentation(template_path)
        
        replacements = {
            'Doctor_1': f'ВРАЧ: {doctor_1}' if doctor_1 else 'Doctor_1',
            'Doctor_2': f'ВРАЧ: {doctor_2}' if doctor_2 else 'Doctor_2',
            'Doctor_3': f'ВРАЧ: {doctor_3}' if doctor_3 else 'Doctor_3',
            'Doctor_4': f'ВРАЧ: {doctor_4}' if doctor_4 else 'Doctor_4',
            'Patient_1': f'ПАЦИЕНТ: {patient_1.upper()}' if patient_1 else 'Patient_1',
            'Patient_2': f'ПАЦИЕНТ: {patient_2.upper()}' if patient_2 else 'Patient_2',
            'Patient_3': f'ПАЦИЕНТ: {patient_3.upper()}' if patient_3 else 'Patient_3',
            'Patient_4': f'ПАЦИЕНТ: {patient_4.upper()}' if patient_4 else 'Patient_4',
            'Дата': f'«{day}» {month} {year} г.'
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            if key in run.text:
                                run.text = run.text.replace(key, value)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_output:
            temp_output_path = temp_output.name
        
        prs.save(temp_output_path)
        
        output_filename = "Бланк Врача на печать.pptx"
        
        def cleanup_temp_file():
            try:
                os.unlink(temp_output_path)
            except OSError:
                pass
        
        media_type = (
            'application/vnd.openxmlformats-officedocument.'
            'presentationml.presentation'
        )
        
        background_tasks = BackgroundTasks()
        background_tasks.add_task(cleanup_temp_file)
        
        return FileResponse(
            path=temp_output_path,
            filename=output_filename,
            media_type=media_type,
            background=background_tasks
        )
        
    except Exception as e:
        status = f"Ошибка обработки файла: {str(e)}"
        response = RedirectResponse(url="/doctor_form", status_code=303)
        encoded_status = base64.b64encode(status.encode('utf-8')).decode('ascii')
        response.set_cookie("doctor_form_status", encoded_status, max_age=10)
        return response
